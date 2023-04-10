import os
import numpy as np
from natsort import natsorted
from pptx import Presentation
from pptx.util import Inches
import datetime

def read_pic():
    path_name = input("请输入所需命名图片文件夹绝对路径：")
    path = os.path.join(path_name)
    img_str = ""
    img_list = os.listdir(path)
    img_list = natsorted(img_list)
    for img_name in img_list:
        if img_name[-4:] == '.jpg':
            img_path = os.path.join(path, img_name)
            img_str += img_path + " "
            img_full_name = img_str.split(" ")
    return img_full_name[:-1]

def insert_pic(slide_num, img_name, left, top, width, height):
    '''插入图片'''
    left, top, width, height = np.array([left, top, width, height])/2.54
    pic = slide_num.shapes.add_picture(image_file=img_name,
                                     left=Inches(left),
                                     top=Inches(top),
                                     width=Inches(width),
                                     height=Inches(height)
                                     )

def read_1_slide():
    img_list = read_pic()
    index = 1
    try:
        size = input("输入位置与大小：")
        size_list = size.split(" ")
        number = int(input("请输入想要从第几页开始："))
        for i, slide in enumerate(prs.slides, start=1):
            if i <= number - 1 :
                pass
            else:
                insert_pic(slide, img_list[index-1], float(size_list[0]), float(size_list[1]), float(size_list[3]), float(size_list[2]))
                index += 1
    except IndexError:
        print("图片处理完毕")
    else:
        print("插入图片成功")
    pass

def read_2_slide():
    img_list = read_pic()
    index = 1
    try:
        size_1 = input("输入第一张位置与大小：")
        size_2 = input("输入第二张位置与大小：")
        size_1_list = size_1.split(" ")
        size_2_list = size_2.split(" ")
        number = int(input("请输入想要从第几页开始："))
        for i, slide in enumerate(prs.slides, start=1):
            if i <= number - 1:
                pass
            else:
                insert_pic(slide, img_list[index - 1], float(size_1_list[0]), float(size_1_list[1]),
                           float(size_1_list[3]), float(size_1_list[2]))
                insert_pic(slide, img_list[index], float(size_2_list[0]), float(size_2_list[1]), float(size_2_list[3]),
                           float(size_2_list[2]))
                index += 2
    except IndexError:
        print("图片处理完毕")
    else:
        print("插入图片成功")
    pass

def read_3_slide():
    img_list = read_pic()
    index = 1
    try:
        size_1 = input("输入第一张位置与大小：")
        size_2 = input("输入第二张位置与大小：")
        size_3 = input("输入第三张位置与大小：")
        size_1_list = size_1.split(" ")
        size_2_list = size_2.split(" ")
        size_3_list = size_3.split(" ")
        number = int(input("请输入想要从第几页开始："))
        for i, slide in enumerate(prs.slides, start=1):
            if i <= number - 1:
                pass
            else:
                insert_pic(slide, img_list[index - 1], float(size_1_list[0]), float(size_1_list[1]),
                           float(size_1_list[3]), float(size_1_list[2]))
                insert_pic(slide, img_list[index], float(size_2_list[0]), float(size_2_list[1]), float(size_2_list[3]),
                           float(size_2_list[2]))
                insert_pic(slide, img_list[index + 1], float(size_3_list[0]), float(size_3_list[1]), float(size_3_list[3]),
                           float(size_3_list[2]))
                index += 3

    except IndexError:
        print("图片处理完毕")
    else:
        print("插入图片成功")
    pass

def read_4_slide():
    img_list = read_pic()
    index = 1
    try:
        size_1 = input("输入第一张位置与大小：")
        size_2 = input("输入第二张位置与大小：")
        size_3 = input("输入第三张位置与大小：")
        size_4 = input("输入第四张位置与大小：")
        size_1_list = size_1.split(" ")
        size_2_list = size_2.split(" ")
        size_3_list = size_3.split(" ")
        size_4_list = size_4.split(" ")
        number = int(input("请输入想要从第几页开始："))
        for i, slide in enumerate(prs.slides, start=1):
            if i <= number - 1:
                pass
            else:
                insert_pic(slide, img_list[index - 1], float(size_1_list[0]), float(size_1_list[1]),float(size_1_list[3]),
                           float(size_1_list[2]))
                insert_pic(slide, img_list[index], float(size_2_list[0]), float(size_2_list[1]), float(size_2_list[3]),
                           float(size_2_list[2]))
                insert_pic(slide, img_list[index + 1], float(size_3_list[0]), float(size_3_list[1]), float(size_3_list[3]),
                           float(size_3_list[2]))
                insert_pic(slide, img_list[index + 2], float(size_4_list[0]), float(size_4_list[1]),float(size_4_list[3]),
                           float(size_4_list[2]))
                index += 4

    except IndexError:
        print("图片处理完毕")
    else:
        print("插入图片成功")
    pass

if __name__ == '__main__':
    prs_path = input("请输入修改PPT的地址及名称：")
    prs = Presentation(prs_path)
    select = input("请输入一张ppt中插入几张图片：")
    save_path = prs_path.split(".")[0]
    now = datetime.datetime.now()
    if select == '1':
        read_1_slide()
        prs.save(save_path + '_' + str(now.hour) + "-" + str(now.minute) + '.pptx')
    elif select == '2':
        read_2_slide()
        prs.save(save_path + '_' + str(now.hour) + "-" + str(now.minute) + '.pptx')
    elif select == '3':
        read_3_slide()
        prs.save(save_path + '_' + str(now.hour) + "-" + str(now.minute) + '.pptx')
    elif select == '4':
        read_4_slide()
        prs.save(save_path + '_' + str(now.hour) + "-" + str(now.minute) + '.pptx')
